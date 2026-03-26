"""
╔══════════════════════════════════════════════════════════════════════╗
║        SMART TEACHER — Constructeur de Cours v2                    ║
║                                                                      ║
║  AMÉLIORATIONS v2 :                                                  ║
║    ✅ Chargement direct depuis courses/dm/ch1..ch7 (sans GPT)        ║
║    ✅ Préservation de la structure PPTX slide par slide              ║
║    ✅ Subject forcé à 'data_mining' pour les cours DM               ║
║    ✅ GPT Structurer avec prompt DM/CS (pas générique)               ║
║    ✅ build_from_dm_folder() : pipeline complet depuis courses/dm/   ║
║    ✅ Détection automatique du niveau (M2 → université)             ║
║    ✅ Concepts DM extraits intelligemment (algorithmes, métriques)   ║
╚══════════════════════════════════════════════════════════════════════╝
"""

import json
import logging
import os
import re
import time
from pathlib import Path
from typing import Optional

from openai import AsyncOpenAI

log = logging.getLogger("SmartTeacher.CourseBuilder")

# Chapitres DM
DM_CHAPTERS = {
    1: "Introduction",
    2: "Data, Dataset, Data Warehouse",
    3: "Exploratory Data Analysis",
    4: "Data Cleaning & Preprocessing",
    5: "Feature Engineering",
    6: "Supervised Machine Learning",
    7: "Unsupervised Machine Learning",
}


# ══════════════════════════════════════════════════════════════════════
#  EXTRACTEUR DE TEXTE PDF/DOCX/PPTX
# ══════════════════════════════════════════════════════════════════════

class TextExtractor:
    """Extrait le texte structuré depuis PDF, DOCX, PPTX, TXT."""

    def extract(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        log.info(f"📄 Extraction : {Path(file_path).name}")
        if ext == ".pdf":
            return self._extract_pdf(file_path)
        elif ext == ".docx":
            return self._extract_docx(file_path)
        elif ext == ".pptx":
            return self._extract_pptx(file_path)
        elif ext in (".txt", ".md"):
            return Path(file_path).read_text(encoding="utf-8", errors="ignore")
        else:
            raise ValueError(f"Format non supporté : {ext}")

    def extract_structured_pptx(self, path: str) -> list[dict]:
        """
        Extrait le PPTX slide par slide.
        Retourne une liste de dicts {slide_idx, title, bullets, content}.
        Préserve la structure des slides pour le SlideSync.
        """
        try:
            from pptx import Presentation
            prs = Presentation(path)
            slides = []
            for i, slide in enumerate(prs.slides):
                title_text   = ""
                bullet_texts = []

                for shape in slide.shapes:
                    if not hasattr(shape, "text") or not shape.text.strip():
                        continue
                    text = shape.text.strip()
                    # Le premier texte grand = titre (souvent placeholder title)
                    if (shape.shape_type == 13 or
                        (hasattr(shape, "placeholder_format") and
                         shape.placeholder_format and
                         shape.placeholder_format.idx == 0)):
                        title_text = text
                    else:
                        bullet_texts.append(text)

                if not title_text and bullet_texts:
                    title_text = bullet_texts.pop(0)

                content = "\n".join(bullet_texts)
                if title_text or content:
                    slides.append({
                        "slide_idx": i + 1,
                        "title":     title_text,
                        "bullets":   bullet_texts,
                        "content":   f"{title_text}\n{content}".strip(),
                    })
            log.info(f"  ✅ {len(slides)} slides extraites")
            return slides
        except ImportError:
            raise ImportError("Installez python-pptx : pip install python-pptx")

    def _extract_pdf(self, path: str) -> str:
        try:
            import pypdf
            text = ""
            with open(path, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    text += (page.extract_text() or "") + "\n\n"
            return text.strip()
        except ImportError:
            pass
        try:
            from pdfminer.high_level import extract_text
            return extract_text(path)
        except ImportError:
            raise ImportError("Installez pypdf : pip install pypdf")

    def _extract_docx(self, path: str) -> str:
        try:
            import docx
            doc = docx.Document(path)
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            raise ImportError("Installez python-docx : pip install python-docx")

    def _extract_pptx(self, path: str) -> str:
        slides = self.extract_structured_pptx(path)
        return "\n\n".join(
            f"[Slide {s['slide_idx']}]\n{s['content']}"
            for s in slides if s['content']
        )


# ══════════════════════════════════════════════════════════════════════
#  STRUCTUREUR GPT — SPÉCIALISÉ DATA MINING
# ══════════════════════════════════════════════════════════════════════

class GPTStructurer:
    """
    Utilise GPT pour transformer un texte brut en cours structuré.
    Spécialisé pour le domaine Data Mining / Informatique.
    """

    # Prompt DM/CS spécialisé
    DM_SYSTEM_PROMPT = """Tu es un expert en Data Mining, Machine Learning et Informatique.
Tu reçois le contenu brut d'un cours universitaire de Data Mining (M2).
Tu dois le transformer en un cours structuré PRÉSENTABLE À VOIX HAUTE par un professeur IA.

RÈGLES IMPORTANTES :
- Chaque section doit être rédigée comme un professeur qui PARLE en cours (pas comme un manuel)
- CONSERVE tous les termes techniques : k-means, SVM, PCA, AUC, Gini, etc.
- Le niveau est Master 2 (M2) — vocabulaire technique avancé autorisé
- Chaque section = environ 2-3 minutes de présentation orale
- Les concepts doivent inclure des algorithmes, métriques, et exemples concrets DM/ML
- NE SIMPLIFIE PAS la terminologie technique

Réponds UNIQUEMENT en JSON valide, sans texte avant ou après.
Format JSON requis :
{
  "title": "Titre du cours",
  "subject": "data_mining",
  "level": "université",
  "language": "en",
  "chapters": [
    {
      "title": "Titre du chapitre",
      "order": 1,
      "sections": [
        {
          "title": "Titre de la section",
          "order": 1,
          "content": "Texte rédigé pour être LU À VOIX HAUTE. Minimum 3-4 phrases naturelles et complètes avec la terminologie DM précise.",
          "duration_s": 120,
          "concepts": [
            {
              "term": "Terme technique DM/ML",
              "definition": "Définition précise et complète",
              "example": "Exemple concret appliqué à DM/ML",
              "type": "definition|algorithm|metric|formula"
            }
          ]
        }
      ]
    }
  ]
}"""

    GENERIC_SYSTEM_PROMPTS = {
        "fr": """Tu es un expert en pédagogie. 
Tu reçois le texte brut d'un cours.
Transforme-le en cours structuré PRÉSENTABLE À VOIX HAUTE.
Réponds UNIQUEMENT en JSON valide avec ce format :
{
  "title": "Titre", "subject": "cs", "level": "université", "language": "fr",
  "chapters": [{"title":"..","order":1,"sections":[{"title":"..","order":1,
  "content":"Texte oral naturel minimum 3 phrases.","duration_s":120,
  "concepts":[{"term":"..","definition":"..","example":"..","type":"definition"}]}]}]
}""",
        "en": """You are a pedagogy expert.
Transform the raw course text into a structured course PRESENTABLE OUT LOUD.
Reply ONLY with valid JSON using this format:
{
  "title": "Title", "subject": "cs", "level": "university", "language": "en",
  "chapters": [{"title":"..","order":1,"sections":[{"title":"..","order":1,
  "content":"Natural oral text minimum 3 sentences.","duration_s":120,
  "concepts":[{"term":"..","definition":"..","example":"..","type":"definition"}]}]}]
}""",
    }

    def __init__(self):
        self.client = AsyncOpenAI(api_key=os.getenv("OPENAI_API_KEY"))

    async def structure(
        self,
        raw_text: str,
        language: str = "en",
        level:    str = "université",
        title:    str = "",
        subject:  str = "data_mining",
        chapter_idx: int | None = None,
    ) -> dict:
        """
        Envoie le texte brut à GPT et récupère le cours structuré.
        Si subject=data_mining → utilise le prompt DM spécialisé.
        """
        log.info(f"🧠 Structuration GPT ({len(raw_text)} chars, subj={subject})…")
        start = time.time()

        # Choisir le prompt
        if subject == "data_mining":
            system = self.DM_SYSTEM_PROMPT
        else:
            lang = language[:2].lower()
            system = self.GENERIC_SYSTEM_PROMPTS.get(lang, self.GENERIC_SYSTEM_PROMPTS["en"])

        # Limiter la taille
        MAX_CHARS = 14000
        if len(raw_text) > MAX_CHARS:
            log.warning(f"⚠️  Texte tronqué : {len(raw_text)} → {MAX_CHARS} chars")
            raw_text = raw_text[:MAX_CHARS] + "\n\n[... texte tronqué ...]"

        # Contexte chapitre DM
        ch_context = ""
        if chapter_idx and chapter_idx in DM_CHAPTERS:
            ch_context = f"Ce contenu correspond au Chapitre {chapter_idx} : {DM_CHAPTERS[chapter_idx]}.\n"

        user_prompt = (
            f"{ch_context}"
            f"Titre suggéré : {title or 'à déterminer'}\n"
            f"Niveau : {level}\nLangue : {language}\n\n"
            f"TEXTE DU COURS :\n{raw_text}\n\n"
            f"Génère le JSON structuré maintenant :"
        )

        try:
            response = await self.client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": system},
                    {"role": "user",   "content": user_prompt},
                ],
                temperature=0.2,
                max_tokens=4000,
                response_format={"type": "json_object"},
            )
            raw_json = response.choices[0].message.content
            data = json.loads(raw_json)

            # Forcer subject=data_mining
            if subject == "data_mining":
                data["subject"] = "data_mining"
                data["level"]   = "université"

            elapsed  = time.time() - start
            chapters = len(data.get("chapters", []))
            sections = sum(len(ch.get("sections", [])) for ch in data.get("chapters", []))
            log.info(f"✅ Structuration : {chapters} chapitres, {sections} sections ({elapsed:.1f}s)")
            return data

        except json.JSONDecodeError as exc:
            log.error(f"❌ JSON invalide : {exc}")
            raise
        except Exception as exc:
            log.error(f"❌ GPT erreur : {exc}")
            raise


# ══════════════════════════════════════════════════════════════════════
#  CONSTRUCTEUR DE COURS PRINCIPAL
# ══════════════════════════════════════════════════════════════════════

class CourseBuilder:
    """
    Pipeline complet : fichier/dossier → cours présentable par l'IA.

    Pour le cours Data Mining :
        builder = CourseBuilder()
        # Charger tous les chapitres ch1..ch7
        courses = await builder.build_dm_course("C:/Users/Admin/.../courses/dm")

    Pour un fichier unique :
        course_data = await builder.build_from_file("ch1.pdf", language="en")
    """

    def __init__(self):
        self.extractor  = TextExtractor()
        self.structurer = GPTStructurer()

    # ── Pipeline Data Mining complet ──────────────────────────────────
    async def build_dm_course(
        self,
        courses_dm_path: str,
        language: str = "en",
    ) -> dict[int, dict]:
        """
        Charge tous les chapitres DM (ch1..ch7) depuis courses/dm/.
        Retourne un dict {chapter_idx: course_data}.

        Structure attendue :
            courses/dm/
              ch1.pdf  (ou ch1/, Chapter_1.pdf, etc.)
              ch2.pdf
              ...
              ch7.pdf
        """
        dm_path = Path(courses_dm_path)
        if not dm_path.exists():
            raise FileNotFoundError(f"Dossier introuvable : {dm_path}")

        log.info(f"\n{'='*60}")
        log.info(f"📚 Construction DM : {dm_path}")
        log.info(f"{'='*60}")

        results: dict[int, dict] = {}

        for ch_idx, ch_title in DM_CHAPTERS.items():
            file_path = self._find_chapter_file(dm_path, ch_idx)
            if not file_path:
                log.warning(f"⚠️  Chapitre {ch_idx} introuvable, ignoré")
                continue

            log.info(f"\n📖 Ch{ch_idx} — {ch_title} : {file_path.name}")
            try:
                course_data = await self._build_chapter(
                    file_path, ch_idx, ch_title, language
                )
                results[ch_idx] = course_data
                log.info(f"  ✅ Ch{ch_idx} structuré")
            except Exception as exc:
                log.error(f"  ❌ Ch{ch_idx} échoué : {exc}")

        log.info(f"\n✅ DM Course prêt : {len(results)}/7 chapitres chargés")
        return results

    async def _build_chapter(
        self,
        file_path: Path,
        chapter_idx: int,
        chapter_title: str,
        language: str,
    ) -> dict:
        """Construit un chapitre DM depuis un fichier."""
        # Extraction spéciale PPTX (slide par slide)
        if file_path.suffix.lower() == ".pptx":
            slides = self.extractor.extract_structured_pptx(str(file_path))
            raw_text = "\n\n".join(
                f"[Slide {s['slide_idx']}]\n{s['content']}"
                for s in slides
            )
        else:
            raw_text = self.extractor.extract(str(file_path))

        if len(raw_text.strip()) < 100:
            raise ValueError(f"Texte trop court : {len(raw_text)} chars")

        course_data = await self.structurer.structure(
            raw_text=raw_text,
            language=language,
            level="université",
            title=f"Chapter {chapter_idx}: {chapter_title}",
            subject="data_mining",
            chapter_idx=chapter_idx,
        )
        course_data["chapter_idx"]   = chapter_idx
        course_data["chapter_title"] = chapter_title
        course_data["file_path"]     = str(file_path)
        return course_data

    def _find_chapter_file(self, dm_path: Path, ch_idx: int) -> Path | None:
        """Trouve le fichier d'un chapitre dans le dossier DM."""
        patterns = [
            f"ch{ch_idx}.pdf", f"ch{ch_idx}.pptx", f"ch{ch_idx}.docx",
            f"CH{ch_idx}.pdf", f"CH{ch_idx}.pptx",
            f"chapter_{ch_idx}.pdf", f"chapter_{ch_idx}.pptx",
            f"Chapter_{ch_idx}.pdf", f"Chapter_{ch_idx}.pptx",
            f"Chapter{ch_idx}.pdf",
            f"DM_ch{ch_idx}.pdf", f"dm_ch{ch_idx}.pdf",
        ]
        # Fichiers directs
        for p in patterns:
            fp = dm_path / p
            if fp.exists():
                return fp

        # Sous-dossier
        for sub in [f"ch{ch_idx}", f"CH{ch_idx}", f"chapter_{ch_idx}", f"Chapter_{ch_idx}"]:
            sub_path = dm_path / sub
            if sub_path.is_dir():
                for ext in ["*.pdf", "*.pptx", "*.docx"]:
                    files = list(sub_path.glob(ext))
                    if files:
                        return files[0]

        # Chercher par numéro dans le nom
        for f in dm_path.iterdir():
            name = f.name.lower()
            if (f"ch{ch_idx}" in name or f"chapter{ch_idx}" in name or
                f"_{ch_idx}." in name or f"{ch_idx}." in name) and \
               f.suffix.lower() in (".pdf", ".pptx", ".docx"):
                return f

        return None

    # ── Pipeline fichier unique ───────────────────────────────────────
    async def build_from_file(
        self,
        file_path: str,
        language:  str = "en",
        level:     str = "université",
        subject:   str = "data_mining",
    ) -> dict:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"Fichier introuvable : {file_path}")

        log.info(f"\n{'='*60}")
        log.info(f"📚 Construction : {path.name}")
        log.info(f"{'='*60}")

        raw_text = self.extractor.extract(file_path)
        log.info(f"   ✅ {len(raw_text)} caractères extraits")

        if len(raw_text.strip()) < 100:
            raise ValueError(f"Texte trop court : {len(raw_text)} chars")

        course_data = await self.structurer.structure(
            raw_text=raw_text,
            language=language,
            level=level,
            title=path.stem.replace("_", " ").replace("-", " ").title(),
            subject=subject,
        )
        course_data["file_path"] = str(file_path)
        self._print_summary(course_data)
        return course_data

    async def build_from_text(
        self,
        text:     str,
        title:    str = "Data Mining Course",
        language: str = "en",
        level:    str = "université",
        subject:  str = "data_mining",
    ) -> dict:
        log.info(f"📝 Construction depuis texte : {len(text)} chars")
        return await self.structurer.structure(
            raw_text=text, language=language, level=level,
            title=title, subject=subject,
        )

    async def save_to_database(self, course_data: dict, db) -> str:
        from database.models import Course, Chapter, Section, Concept
        log.info("💾 Sauvegarde PostgreSQL…")

        course = Course(
            title=course_data.get("title", "Data Mining"),
            subject=course_data.get("subject", "data_mining"),
            language=course_data.get("language", "en"),
            level=course_data.get("level", "université"),
            description=course_data.get("description", ""),
            file_path=course_data.get("file_path", ""),
        )
        db.add(course)
        await db.flush()

        for ch_data in course_data.get("chapters", []):
            chapter = Chapter(
                course_id=course.id,
                title=ch_data["title"],
                order=ch_data.get("order", 0),
                summary=ch_data.get("summary", ""),
            )
            db.add(chapter)
            await db.flush()

            for sec_data in ch_data.get("sections", []):
                section = Section(
                    chapter_id=chapter.id,
                    title=sec_data["title"],
                    order=sec_data.get("order", 0),
                    content=sec_data.get("content", ""),
                    duration_s=sec_data.get("duration_s", 120),
                )
                db.add(section)
                await db.flush()

                for c_data in sec_data.get("concepts", []):
                    concept = Concept(
                        section_id=section.id,
                        term=c_data.get("term", ""),
                        definition=c_data.get("definition", ""),
                        example=c_data.get("example", ""),
                        concept_type=c_data.get("type", "definition"),
                    )
                    db.add(concept)

        await db.commit()
        course_id = str(course.id)
        log.info(f"✅ Cours sauvegardé : ID={course_id}")
        return course_id

    def _print_summary(self, data: dict) -> None:
        chapters = data.get("chapters", [])
        total_sections = sum(len(ch.get("sections", [])) for ch in chapters)
        total_concepts = sum(
            len(sec.get("concepts", []))
            for ch in chapters
            for sec in ch.get("sections", [])
        )
        total_duration = sum(
            sec.get("duration_s", 120)
            for ch in chapters
            for sec in ch.get("sections", [])
        )
        log.info(f"\n{'='*60}")
        log.info(f"✅ COURS : {data.get('title')}")
        log.info(f"   Matière    : {data.get('subject')}")
        log.info(f"   Niveau     : {data.get('level')}")
        log.info(f"   Langue     : {data.get('language')}")
        log.info(f"   Chapitres  : {len(chapters)}")
        log.info(f"   Sections   : {total_sections}")
        log.info(f"   Concepts   : {total_concepts}")
        log.info(f"   Durée est. : {total_duration // 60} min")
        log.info(f"{'='*60}\n")
        for i, ch in enumerate(chapters):
            log.info(f"  📖 Ch{i+1} : {ch['title']}")
            for j, sec in enumerate(ch.get("sections", [])):
                log.info(f"     └─ §{j+1} : {sec['title']} ({sec.get('duration_s',120)}s)")