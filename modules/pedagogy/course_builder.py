"""
╔══════════════════════════════════════════════════════════════════════╗
║        SMART TEACHER — Constructeur de Cours v3                    ║
║                                                                      ║
║  AMÉLIORATIONS v3 :                                                  ║
║    ✅ Structure hiérarchique : Domain → Courses → Chapters          ║
║    ✅ Support multi-domaines                                     ║
║    ✅ Chargement depuis courses/{domain}/{course}/                  ║
║    ✅ Préservation de la structure PPTX slide par slide              ║
║    ✅ Subject automatique selon le domaine & cours                  ║
║    ✅ Pipeline complet : build_course_chapters()                    ║
║    ✅ Détection automatique du niveau (université)                  ║
║    ✅ Concepts extraits intelligemment                              ║
╚══════════════════════════════════════════════════════════════════════╝
"""

import json
import logging
import os
import re
import time
import warnings
from pathlib import Path
from typing import Optional

# Import configuration domaines & cours
import sys
sys.path.insert(0, str(Path(__file__).parent.parent))
from domains_config import DEFAULT_DOMAIN, DEFAULT_COURSE, get_chapters, get_courses, get_domains

log = logging.getLogger("SmartTeacher.CourseBuilder")


# ══════════════════════════════════════════════════════════════════════
#  EXTRACTEUR DE TEXTE PDF/DOCX/PPTX
# ══════════════════════════════════════════════════════════════════════

class TextExtractor:
    """Extrait le texte structuré depuis PDF, DOCX, PPTX, TXT."""

    def extract(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        log.info(f"📄 Extraction : {Path(file_path).name}")
        if ext == ".pdf":
            return self._extract_pdf(file_path)
        elif ext == ".docx":
            return self._extract_docx(file_path)
        elif ext == ".pptx":
            return self._extract_pptx(file_path)
        elif ext in (".txt", ".md"):
            return Path(file_path).read_text(encoding="utf-8", errors="ignore")
        else:
            raise ValueError(f"Format non supporté : {ext}")

    def extract_structured_pptx(self, path: str) -> list[dict]:
        """
        Extrait le PPTX slide par slide.
        Retourne une liste de dicts {slide_idx, title, bullets, content}.
        Préserve la structure des slides pour le SlideSync.
        """
        try:
            from pptx import Presentation
            prs = Presentation(path)
            slides = []
            for i, slide in enumerate(prs.slides):
                title_text   = ""
                bullet_texts = []

                for shape in slide.shapes:
                    if not hasattr(shape, "text") or not shape.text.strip():
                        continue
                    text = shape.text.strip()
                    is_title_placeholder = shape.shape_type == 13
                    if not is_title_placeholder:
                        try:
                            placeholder_format = shape.placeholder_format
                        except (AttributeError, ValueError):
                            placeholder_format = None
                        is_title_placeholder = bool(placeholder_format and placeholder_format.idx == 0)

                    # Le premier texte grand = titre (souvent placeholder title)
                    if is_title_placeholder:
                        title_text = text
                    else:
                        bullet_texts.append(text)

                if not title_text and bullet_texts:
                    title_text = bullet_texts.pop(0)

                content = "\n".join(bullet_texts)
                if title_text or content:
                    slides.append({
                        "slide_idx": i + 1,
                        "title":     title_text,
                        "bullets":   bullet_texts,
                        "content":   f"{title_text}\n{content}".strip(),
                    })
            log.info(f"  ✅ {len(slides)} slides extraites")
            return slides
        except ImportError:
            raise ImportError("Installez python-pptx : pip install python-pptx")

    def _extract_pdf(self, path: str) -> str:
        try:
            import pypdf
            text = ""
            with open(path, "rb") as f:
                with warnings.catch_warnings():
                    warnings.filterwarnings("ignore", module=r"pypdf\.generic\._base")
                    reader = pypdf.PdfReader(f, strict=False)
                    for page in reader.pages:
                        text += (page.extract_text() or "") + "\n\n"
            return text.strip()
        except ImportError:
            pass
        try:
            from pdfminer.high_level import extract_text
            return extract_text(path)
        except ImportError:
            raise ImportError("Installez pypdf : pip install pypdf")

    def _extract_docx(self, path: str) -> str:
        try:
            import docx
            doc = docx.Document(path)
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            raise ImportError("Installez python-docx : pip install python-docx")

    def _extract_pptx(self, path: str) -> str:
        slides = self.extract_structured_pptx(path)
        return "\n\n".join(
            f"[Slide {s['slide_idx']}]\n{s['content']}"
            for s in slides if s['content']
        )


# ══════════════════════════════════════════════════════════════════════
#  STRUCTUREUR LOCAL — TESSERACT OCR (sans API OpenAI)
# ══════════════════════════════════════════════════════════════════════

class LocalStructurer:
    """
    ✅ Structure le texte SANS modifier le contenu original.
    ✅ Utilise Tesseract-OCR local (gratuit, rapide).
    ✅ Génère les PNG pour l'affichage.
    ✅ Indexe le texte brut directement sans appel OpenAI.
    """

    def __init__(self):
        """Initialise Tesseract (déjà dans le Dockerfile)"""
        try:
            import pytesseract
            from pdf2image import convert_from_path
            self.pytesseract = pytesseract
            self.convert_from_path = convert_from_path
            self.available_ocr_langs = set()
            try:
                self.available_ocr_langs = set(pytesseract.get_languages(config="") or [])
            except Exception:
                self.available_ocr_langs = set()
            log.info("✅ Tesseract-OCR en local activé")
        except ImportError as e:
            log.info(f"ℹ️ Tesseract non disponible: {e}")
            self.pytesseract = None
            self.convert_from_path = None
            self.available_ocr_langs = set()

    def _select_ocr_language(self, requested_lang: str) -> str:
        """Choisit une langue Tesseract réellement installée pour éviter le fallback English par défaut."""
        if not self.pytesseract or not self.available_ocr_langs:
            return ""

        requested_parts = [part.strip() for part in requested_lang.split("+") if part.strip()]
        selected = [part for part in requested_parts if part in self.available_ocr_langs]
        if selected:
            return "+".join(selected)

        for fallback in ("fra", "ara", "eng"):
            if fallback in self.available_ocr_langs:
                return fallback

        return next(iter(sorted(self.available_ocr_langs)), "")

    def _pdf_to_images(self, pdf_path: str, output_dir: str) -> list[str]:
        """Convertit PDF en images PNG."""
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)

        for old_png in output_path.glob("page_*.png"):
            try:
                old_png.unlink()
            except Exception:
                pass

        if self.convert_from_path:
            try:
                images = self.convert_from_path(pdf_path, dpi=150)
                paths = []
                for i, img in enumerate(images, 1):
                    png_path = output_path / f"page_{i:03d}.png"
                    img.save(str(png_path), "PNG")
                    paths.append(str(png_path))
                if paths:
                    log.info(f"  ✅ {len(paths)} PNG générées")
                    return paths
                log.warning("  ⚠️ pdf2image n'a produit aucune image, fallback pypdfium2…")
            except Exception as e:
                log.info(f"  ℹ️ Conversion PDF→PNG via pdf2image échouée: {e}")

        try:
            import pypdfium2 as pdfium

            pdf = pdfium.PdfDocument(pdf_path)
            paths = []
            try:
                for i in range(len(pdf)):
                    page = pdf[i]
                    bitmap = page.render(scale=2.0)
                    image = bitmap.to_pil()
                    png_path = output_path / f"page_{i + 1:03d}.png"
                    image.save(str(png_path), "PNG")
                    paths.append(str(png_path))
            finally:
                close = getattr(pdf, "close", None)
                if callable(close):
                    close()

            if paths:
                log.info(f"  ✅ {len(paths)} PNG générées via pypdfium2")
            return paths
        except ImportError:
            log.info("  ℹ️ pypdfium2 non disponible pour convertir le PDF en PNG")
        except Exception as e:
            log.info(f"  ℹ️ Conversion PDF→PNG via pypdfium2 échouée: {e}")
            return []

    def _extract_text_with_ocr(self, pdf_path: str, lang: str = "fra+ara+eng") -> str:
        """Extrait le texte brut avec Tesseract-OCR."""
        if not self.pytesseract or not self.convert_from_path:
            return ""
        try:
            resolved_lang = self._select_ocr_language(lang)
            if not resolved_lang:
                log.info("ℹ️ OCR ignoré: aucun pack de langue Tesseract disponible")
                return ""

            images = self.convert_from_path(pdf_path, dpi=150)
            full_text = ""
            for img in images:
                # Utiliser Tesseract pour extraire le texte
                text = self.pytesseract.image_to_string(img, lang=resolved_lang)
                full_text += text + "\n\n"
            return full_text.strip()
        except Exception as e:
            log.info(f"ℹ️ OCR échouée: {e}")
            return ""

    async def structure(
        self,
        raw_text: str,
        language: str = "en",
        level: str = "université",
        title: str = "",
        subject: str = DEFAULT_COURSE,
        chapter_idx: int | None = None,
    ) -> dict:
        """
        Structure le texte SANS le modifier (100% local).
        
        Retourne une structure simple pour l'indexation RAG.
        Le contenu est gardé intégralement.
        """
        log.info(f"📚 Structuration locale ({len(raw_text)} chars, tesseract)…")
        
        # Diviser le texte en sections (~ 500 mots chacune)
        sentences = [s.strip() for s in raw_text.replace("\n\n", ". ").split(". ") if s.strip()]
        
        sections = []
        current_section = []
        word_count = 0
        
        for sentence in sentences:
            words = sentence.split()
            current_section.append(sentence)
            word_count += len(words)
            
            if word_count >= 150:  # ~1 minute orale
                content = ". ".join(current_section)
                if content.strip():
                    sections.append({
                        "title": content.split("\n")[0][:60] + "...",
                        "order": len(sections) + 1,
                        "content": content,
                        "duration_s": max(60, word_count // 3),  # ~3 mots/seconde
                        "concepts": self._extract_concepts(content, subject),
                    })
                current_section = []
                word_count = 0
        
        # Dernière section
        if current_section:
            content = ". ".join(current_section)
            sections.append({
                "title": content.split("\n")[0][:60] + "...",
                "order": len(sections) + 1,
                "content": content,
                "duration_s": max(60, word_count // 3),
                "concepts": self._extract_concepts(content, subject),
            })
        
        return {
            "title": title or "Cours importé",
            "subject": subject,
            "level": level,
            "language": language,
            "chapters": [
                {
                    "title": f"Chapitre {chapter_idx or 1}",
                    "order": chapter_idx or 1,
                    "sections": sections,
                }
            ],
        }

    def _extract_concepts(self, text: str, subject: str = DEFAULT_COURSE) -> list[dict]:
        """Extrait les concepts clés (simple regex, pas d'IA)."""
        concepts = []

        keywords = {
            "general": [
                "concept", "definition", "example", "method", "process",
                "analysis", "model", "result", "variable", "data",
            ],
            "computer_science": [
                "algorithm", "complexity", "recursion", "tree", "graph",
                "database", "network", "protocol", "memory", "function",
            ],
            "mathematics": [
                "theorem", "proof", "equation", "integral", "derivative",
                "matrix", "vector", "probability", "statistics", "formula",
            ],
            "science": [
                "experiment", "measurement", "hypothesis", "observation",
                "analysis", "result", "variable", "sample", "process",
            ],
            "languages": [
                "grammar", "syntax", "vocabulary", "reading", "writing",
                "translation", "text", "sentence", "meaning", "context",
            ],
        }

        key_list = keywords.get(subject, keywords["general"])
        
        for keyword in key_list:
            if keyword.lower() in text.lower():
                concepts.append({
                    "term": keyword,
                    "definition": f"Terme trouvé dans le contenu",
                    "example": "Voir le texte du cours",
                    "type": "keyword",
                })
        
        return concepts[:5]  # Max 5 concepts par section


# ══════════════════════════════════════════════════════════════════════
#  CONSTRUCTEUR DE COURS PRINCIPAL
# ══════════════════════════════════════════════════════════════════════

class CourseBuilder:
    """
    Pipeline complet : fichier/dossier → cours présentable par l'IA.

    Exemple d'usage :
        builder = CourseBuilder()
        course_data = await builder.build_from_file("mon_fichier.pdf", language="en")
    """

    def __init__(self):
        self.extractor  = TextExtractor()
        self.structurer = LocalStructurer()

    @staticmethod
    def _course_slug(value: str | None, fallback: str = DEFAULT_COURSE, domain: str | None = None) -> str:
        candidate = (value or fallback).strip()

        if domain:
            try:
                for existing_course in get_courses(domain):
                    if candidate == existing_course or candidate.lower() == existing_course.lower():
                        return existing_course
            except Exception:
                pass

        candidate = candidate.lower()
        candidate = re.sub(r"[^a-z0-9]+", "_", candidate)
        candidate = candidate.strip("_")
        return candidate or fallback

    @staticmethod
    def _looks_like_chapter(value: str | None) -> bool:
        if not value:
            return False

        normalized = value.strip().lower().replace("\\", "/").split("/")[-1]
        return bool(
            re.search(r"(?:chapter|chapitre|chap)\s*[_\-\s]*\d+", normalized)
            or re.fullmatch(r"chapter[_\-\s]*\d+", normalized)
            or re.fullmatch(r"chapitre[_\-\s]*\d+", normalized)
            or re.fullmatch(r"ch\d+", normalized)
        )

    @staticmethod
    def _chapter_number_from_value(value: str | None) -> int | None:
        if not value:
            return None

        normalized = value.strip().lower().replace("\\", "/").split("/")[-1]
        patterns = [
            r"(?:chapter|chapitre|chap)\s*[_\-\s]*(\d+)",
            r"\bchapter[_\-\s]*(\d+)\b",
            r"\bchapitre[_\-\s]*(\d+)\b",
            r"\bchap[_\-\s]*(\d+)\b",
            r"\bch[_\-\s]*(\d+)\b",
        ]

        for pattern in patterns:
            match = re.search(pattern, normalized)
            if match:
                try:
                    return max(1, int(match.group(1)))
                except (TypeError, ValueError):
                    continue

        return None

    @staticmethod
    def _chapter_slug(value: str | None, fallback: str = "chapter_1") -> str:
        candidate = (value or fallback).strip().replace("\\", "/").split("/")[-1]
        candidate = candidate.lower()
        candidate = re.sub(r"[^a-z0-9]+", "_", candidate)
        candidate = candidate.strip("_")
        return candidate or fallback

    @staticmethod
    def _display_label(value: str | None, fallback: str) -> str:
        candidate = (value or "").replace("_", " ").replace("-", " ").strip()
        return candidate.title() if candidate else fallback

    def _infer_course_title(self, raw_text: str, course_slug: str, fallback_title: str) -> str:
        """Infer a real course title from the detected course slug or the document text."""
        slug = (course_slug or "").strip()
        if slug and slug.lower() not in {DEFAULT_COURSE, "generic"} and not self._looks_like_chapter(slug):
            return self._display_label(slug, fallback_title)

        for line in (raw_text or "").splitlines():
            candidate = line.strip()
            if not candidate:
                continue
            if len(candidate) > 120 or len(candidate.split()) > 10:
                continue
            if self._looks_like_chapter(candidate):
                continue
            if not any(char.isalpha() for char in candidate):
                continue
            normalized = candidate.lower().strip()
            if normalized in {"outline", "table of contents", "contents", "introduction", "summary"}:
                continue
            return self._display_label(candidate, candidate)

        return self._display_label(fallback_title, "Cours importé")

    def infer_upload_context(
        self,
        upload_name: str | None,
        fallback_domain: str = DEFAULT_DOMAIN,
        fallback_course: str | None = None,
        fallback_chapter: str = "chapter_1",
    ) -> tuple[str, str, str]:
        """Déduit domaine / cours / chapitre depuis le chemin uploadé."""
        raw_path = (upload_name or "").replace("\\", "/").strip()
        parts = [part for part in Path(raw_path).parts if part not in ("", ".", "..")]
        if parts and len(parts[0]) == 2 and parts[0][1] == ":":
            parts = parts[1:]

        folders = parts[:-1]
        file_stem = Path(parts[-1]).stem if parts else ""
        known_domains = {name.lower(): name for name in get_domains()}

        domain = fallback_domain or DEFAULT_DOMAIN
        course = fallback_course if fallback_course and not self._looks_like_chapter(fallback_course) else None
        if course is None and file_stem and not self._looks_like_chapter(file_stem):
            course = file_stem
        if course is None:
            course = DEFAULT_COURSE
        chapter = fallback_chapter or "chapter_1"
        chapter_from_folder = False

        if folders:
            first_folder = folders[0]
            matched_domain = known_domains.get(first_folder.lower())
            if matched_domain:
                domain = matched_domain

            chapter_idx = next(
                (idx for idx in range(len(folders) - 1, -1, -1) if self._looks_like_chapter(folders[idx])),
                None,
            )
            if chapter_idx is not None:
                chapter = folders[chapter_idx]
                chapter_from_folder = True
                if chapter_idx >= 1:
                    course = folders[chapter_idx - 1]
                if chapter_idx >= 2:
                    domain = folders[chapter_idx - 2]
            elif len(folders) >= 3:
                domain, course, chapter = folders[-3], folders[-2], folders[-1]
                chapter_from_folder = True
            elif len(folders) == 2:
                first, second = folders
                if self._looks_like_chapter(second):
                    if not matched_domain:
                        course = first
                    chapter = second
                    chapter_from_folder = True
                else:
                    domain, course = first, second
            elif len(folders) == 1:
                if matched_domain:
                    domain = matched_domain
                elif self._looks_like_chapter(folders[0]):
                    chapter = folders[0]
                    chapter_from_folder = True
                else:
                    course = folders[0]

        if file_stem and self._looks_like_chapter(file_stem) and not chapter_from_folder:
            chapter = file_stem

        if not domain:
            domain = fallback_domain or DEFAULT_DOMAIN

        if course and domain.lower() in {DEFAULT_DOMAIN.lower(), "general"}:
            try:
                course_key = self._course_slug(course, fallback=course, domain=None).lower()
                for known_domain in get_domains():
                    if known_domain.lower() == domain.lower():
                        continue
                    try:
                        existing_courses = get_courses(known_domain)
                    except Exception:
                        continue

                    if any(self._course_slug(existing_course, fallback=existing_course, domain=None).lower() == course_key for existing_course in existing_courses):
                        domain = known_domain
                        break
            except Exception:
                pass

        if not course or self._looks_like_chapter(course):
            if fallback_course and not self._looks_like_chapter(fallback_course):
                course = fallback_course
            elif file_stem and not self._looks_like_chapter(file_stem):
                course = file_stem
            else:
                course = DEFAULT_COURSE

        return (
            domain,
            self._course_slug(course, fallback_course or course or DEFAULT_COURSE, domain=domain),
            self._chapter_slug(chapter, fallback_chapter),
        )

    # ── Pipeline de cours complet ─────────────────────────────────────
    async def build_course_chapters(
        self,
        domain: str = DEFAULT_DOMAIN,
        course: str = DEFAULT_COURSE,
        language: str = "en",
        auto_detect: bool = False,
        sample_file: str | None = None,
    ) -> dict[int, dict]:
        """
        Charge tous les chapitres d'un cours depuis courses/{domain}/{course}/.
        Retourne un dict {chapter_idx: course_data}.

        Args:
            domain: Nom du domaine (ex: "informatique")
            course: Nom du cours (ex: "mathematiques")
            language: Langue (ex: "fr", "en")
            auto_detect: Si True, détecte automatiquement le domaine/cours depuis sample_file
            sample_file: Chemin vers un fichier PDF pour auto-détection

        Structure attendue :
            courses/informatique/mathematiques/
              Chapter 1.pdf
              Chapter 2.pdf
              ...
              Chapter N.pdf

        Exemples :
            # Spécifique
            chapters = await builder.build_course_chapters("informatique", "mathematiques", "fr")

            # Auto-détection depuis un PDF
            chapters = await builder.build_course_chapters(
                auto_detect=True,
                sample_file="Chapter 1.pdf",
                language="fr"
            )
        """
        if auto_detect and sample_file:
            log.info(f"🔍 Auto-détection du domaine/cours depuis : {sample_file}")
            try:
                from domains_config import auto_detect_course
                domain, course = auto_detect_course(sample_file)
                log.info(f"✅ Détecté : {domain} / {course}")
            except Exception as e:
                log.info(f"ℹ️  Auto-détection échouée : {e}. Utilisation valeurs par défaut.")
                domain = DEFAULT_DOMAIN
                course = DEFAULT_COURSE

        available_courses = get_courses(domain)
        if course not in available_courses:
            raise ValueError(
                f"Cours '{course}' introuvable dans '{domain}'. "
                f"Disponibles : {available_courses}"
            )

        course_path = Path("courses") / domain / course
        if not course_path.exists():
            raise FileNotFoundError(f"Dossier introuvable : {course_path}")

        log.info(f"\n{'='*60}")
        log.info(f"📚 Construction {domain.upper()} / {course.upper()}")
        log.info(f"📂 Chemin : {course_path}")
        log.info(f"{'='*60}")

        results: dict[int, dict] = {}
        chapters = get_chapters(domain, course)

        for ch_idx, ch_title in chapters.items():
            file_path = self._find_chapter_file(course_path, ch_idx)
            if not file_path:
                log.info(f"ℹ️  Chapitre {ch_idx} introuvable, ignoré")
                continue

            log.info(f"\n📖 Ch{ch_idx} — {ch_title} : {file_path.name}")
            try:
                course_data = await self._build_chapter(
                    file_path,
                    ch_idx,
                    ch_title,
                    language,
                    domain=domain,
                    subject=course,
                )
                results[ch_idx] = course_data
                log.info(f"  ✅ Ch{ch_idx} structuré")
            except Exception as exc:
                log.error(f"  ❌ Ch{ch_idx} échoué : {exc}")

        log.info(f"\n✅ Cours {course} prêt : {len(results)}/{len(chapters)} chapitres chargés")
        return results

    def _infer_context_from_file(
        self,
        file_path: str,
        fallback_domain: str = DEFAULT_DOMAIN,
        fallback_course: str | None = None,
        fallback_chapter: str = "chapter_1",
    ) -> tuple[str, str, str]:
        """Infère la hiérarchie à partir du chemin local sauvegardé."""
        return self.infer_upload_context(
            file_path,
            fallback_domain=fallback_domain,
            fallback_course=fallback_course,
            fallback_chapter=fallback_chapter,
        )

    async def _build_chapter(
        self,
        file_path: Path,
        chapter_idx: int,
        chapter_title: str,
        language: str,
        domain: str = DEFAULT_DOMAIN,
        subject: str = DEFAULT_COURSE,
        chapter: str | None = None,
    ) -> dict:
        """Construit un chapitre depuis un fichier."""
        
        # 🎨 Générer les PNG (images visuelles du cours)
        chapter_slug = self._chapter_slug(chapter or f"chapter_{chapter_idx}", f"chapter_{chapter_idx}")
        course_dir = Path("media/slides") / domain / subject
        chapter_dir = course_dir / chapter_slug
        chapter_dir.mkdir(parents=True, exist_ok=True)
        
        png_paths = self.structurer._pdf_to_images(
            str(file_path),
            str(chapter_dir)
        )
        log.info(f"  📸 {len(png_paths)} PNG générées → {chapter_dir}")
        
        # Extraction spéciale PPTX (slide par slide)
        if file_path.suffix.lower() == ".pptx":
            slides = self.extractor.extract_structured_pptx(str(file_path))
            raw_text = "\n\n".join(
                f"[Slide {s['slide_idx']}]\n{s['content']}"
                for s in slides
            )
        else:
            raw_text = self.extractor.extract(str(file_path))

        if len(raw_text.strip()) < 100:
            raise ValueError(f"Texte trop court : {len(raw_text)} chars")

        course_data = await self.structurer.structure(
            raw_text=raw_text,
            language=language,
            level="université",
            title=f"Chapter {chapter_idx}: {chapter_title}",
            subject=subject,
            chapter_idx=chapter_idx,
        )
        course_data["chapter_idx"]   = chapter_idx
        course_data["chapter_title"] = chapter_title
        course_data["chapter_slug"]  = chapter_slug
        course_data["file_path"]     = str(file_path)
        course_data["slides"]        = [f"/media/slides/{domain}/{subject}/{chapter_slug}/page_{i+1:03d}.png" for i in range(len(png_paths))]
        return course_data

    def _find_chapter_file(self, course_path: Path, ch_idx: int) -> Path | None:
        """Trouve le fichier d'un chapitre dans le dossier du cours."""
        patterns = [
            f"ch{ch_idx}.pdf", f"ch{ch_idx}.pptx", f"ch{ch_idx}.docx",
            f"CH{ch_idx}.pdf", f"CH{ch_idx}.pptx",
            f"chapter_{ch_idx}.pdf", f"chapter_{ch_idx}.pptx",
            f"Chapter_{ch_idx}.pdf", f"Chapter_{ch_idx}.pptx",
            f"Chapter{ch_idx}.pdf",
        ]
        # Fichiers directs
        for p in patterns:
            fp = course_path / p
            if fp.exists():
                return fp

        # Sous-dossier
        for sub in [f"ch{ch_idx}", f"CH{ch_idx}", f"chapter_{ch_idx}", f"Chapter_{ch_idx}"]:
            sub_path = course_path / sub
            if sub_path.is_dir():
                for ext in ["*.pdf", "*.pptx", "*.docx"]:
                    files = list(sub_path.glob(ext))
                    if files:
                        return files[0]

        # Chercher par numéro dans le nom
        for f in course_path.iterdir():
            name = f.name.lower()
            if (f"ch{ch_idx}" in name or f"chapter{ch_idx}" in name or
                f"_{ch_idx}." in name or f"{ch_idx}." in name) and \
               f.suffix.lower() in (".pdf", ".pptx", ".docx"):
                return f

        return None

    # ── Pipeline fichier unique ───────────────────────────────────────
    async def build_from_file(
        self,
        file_path: str,
        language:  str = "en",
        level:     str = "université",
        subject:   str = "",
        domain:    str = DEFAULT_DOMAIN,  # 🎯 Domaine (général, informatique, etc.)
        chapter:   str = "chapter_1",
    ) -> dict:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"Fichier introuvable : {file_path}")

        log.info(f"\n{'='*60}")
        log.info(f"📚 Construction : {path.name}")
        log.info(f"{'='*60}")

        inferred_domain, course_slug, chapter_slug = self._infer_context_from_file(
            str(path),
            fallback_domain=domain,
            fallback_course=subject or None,
            fallback_chapter=chapter,
        )
        domain = inferred_domain or domain

        # 🎨 Générer les PNG (images visuelles du cours)
        course_dir = Path("media/slides") / domain / course_slug
        chapter_dir = course_dir / chapter_slug
        chapter_dir.mkdir(parents=True, exist_ok=True)
        
        png_paths = self.structurer._pdf_to_images(
            file_path,
            str(chapter_dir)
        )
        log.info(f"  📸 {len(png_paths)} PNG générées → {chapter_dir}")

        raw_text = self.extractor.extract(file_path)
        log.info(f"   ✅ {len(raw_text)} caractères extraits")

        if len(raw_text.strip()) < 100:
            raise ValueError(f"Texte trop court : {len(raw_text)} chars")

        course_title = self._infer_course_title(raw_text, course_slug, path.stem)

        course_data = await self.structurer.structure(
            raw_text=raw_text,
            language=language,
            level=level,
            title=course_title,
            subject=course_slug,
        )
        chapter_order = self._chapter_number_from_value(chapter_slug) or 1
        course_data["file_path"] = str(file_path)
        course_data["subject"] = course_slug
        course_data["title"] = course_title
        course_data["chapter_slug"] = chapter_slug
        course_data["chapters"][0]["title"] = self._display_label(chapter_slug, f"Chapter {chapter_order}")
        course_data["chapters"][0]["order"] = chapter_order
        course_data["chapter_order"] = chapter_order
        course_data["slides"] = [f"/media/slides/{domain}/{course_slug}/{chapter_slug}/page_{i+1:03d}.png" for i in range(len(png_paths))]
        self._print_summary(course_data)
        return course_data

    async def build_from_file_direct(
        self,
        file_path: str,
        language: str = "en",
        level: str = "université",
        subject: str = "",
        domain: str = DEFAULT_DOMAIN,
        chapter: str = "chapter_1",
        course_title_hint: str | None = None,
    ) -> dict:
        """
        Construit un cours directement depuis le fichier, sans génération IA.
        - PPTX : 1 slide = 1 section (structure conservée)
        - PDF  : 1 page  = 1 section (ordre exact conservé)
        - DOCX/TXT/MD : sections extraites par heuristique sur les titres
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"Fichier introuvable : {file_path}")

        inferred_domain, course_slug, chapter_slug = self._infer_context_from_file(
            str(path),
            fallback_domain=domain,
            fallback_course=subject or None,
            fallback_chapter=chapter,
        )
        domain = inferred_domain or domain
        chapter_title = self._display_label(chapter_slug, "Chapter 1")
        course_title_hint = (course_title_hint or "").strip()

        if course_slug.lower() in {DEFAULT_COURSE, "generic"} and course_title_hint:
            course_slug = self._course_slug(course_title_hint, fallback=course_slug, domain=domain)

        log.info(f"\n{'='*60}")
        log.info(f"📚 Construction directe (sans IA) : {path.name}")
        log.info(f"{'='*60}")

        sections: list[dict] = []
        ext = path.suffix.lower()
        slides: list[str] = []
        title_source_text = ""

        if ext == ".pdf":
            course_dir = Path("media/slides") / domain / course_slug
            chapter_dir = course_dir / chapter_slug
            chapter_dir.mkdir(parents=True, exist_ok=True)

            png_paths = self.structurer._pdf_to_images(str(path), str(chapter_dir))
            slides = [
                f"/media/slides/{domain}/{course_slug}/{chapter_slug}/page_{i+1:03d}.png"
                for i in range(len(png_paths))
            ]

            pages = self._extract_pdf_pages(str(path))
            title_source_text = "\n\n".join(page_text for _, page_text in pages if page_text)
            for page_idx, page_text in pages:
                content = (page_text or "").strip()
                if not content:
                    continue

                title = self._infer_page_title(content, page_idx)
                sections.append({
                    "title": title,
                    "order": page_idx,
                    "page_index": page_idx,
                    "content": content,
                    "duration_s": self._estimate_duration(content),
                    "concepts": [],
                    "image_url": slides[page_idx - 1] if page_idx - 1 < len(slides) else "",
                })
        elif ext == ".pptx":
            slides_data = self.extractor.extract_structured_pptx(str(path))
            title_source_text = "\n\n".join((s.get("content") or "") for s in slides_data if s.get("content"))
            for i, s in enumerate(slides_data, start=1):
                title = (s.get("title") or f"Slide {i}").strip()
                content = (s.get("content") or "").strip()
                if not content:
                    continue
                sections.append({
                    "title": title,
                    "order": i,
                    "page_index": i,
                    "content": content,
                    "duration_s": self._estimate_duration(content),
                    "concepts": [],
                    "image_url": "",
                })
        else:
            raw_text = self.extractor.extract(str(path)).strip()
            title_source_text = raw_text
            if len(raw_text) < 30:
                raise ValueError(f"Texte trop court : {len(raw_text)} chars")

            chunks = self._split_text_into_sections(raw_text)
            for i, chunk in enumerate(chunks, start=1):
                title, content = chunk
                if not content.strip():
                    continue
                sections.append({
                    "title": title or f"Section {i}",
                    "order": i,
                    "page_index": i,
                    "content": content.strip(),
                    "duration_s": self._estimate_duration(content),
                    "concepts": [],
                    "image_url": "",
                })

        if not sections:
            raise ValueError("Impossible d'extraire des sections depuis ce fichier")

        course_title = self._display_label(course_title_hint, path.stem) if course_title_hint else self._infer_course_title(title_source_text, course_slug, path.stem)

        course_data = {
            "title": course_title,
            "subject": course_slug,
            "language": language,
            "level": level,
            "description": "Cours importé directement (sans génération IA)",
            "chapter_order": self._chapter_number_from_value(chapter_slug) or 1,
            "chapters": [
                {
                    "title": chapter_title,
                    "order": self._chapter_number_from_value(chapter_slug) or 1,
                    "summary": "Import direct du document original",
                    "sections": sections,
                }
            ],
            "file_path": str(path),
            "slides": slides,
            "chapter_slug": chapter_slug,
        }

        self._print_summary(course_data)
        return course_data

    def _extract_pdf_pages(self, path: str) -> list[tuple[int, str]]:
        """Extrait le texte page par page pour conserver l'ordre exact du PDF."""
        try:
            import pypdf
            out: list[tuple[int, str]] = []
            with open(path, "rb") as f:
                with warnings.catch_warnings():
                    warnings.filterwarnings("ignore", module=r"pypdf\.generic\._base")
                    reader = pypdf.PdfReader(f, strict=False)
                    for i, page in enumerate(reader.pages, start=1):
                        out.append((i, (page.extract_text() or "").strip()))
            return out
        except ImportError as exc:
            raise ImportError("Installez pypdf : pip install pypdf") from exc

    async def build_from_text(
        self,
        text:     str,
        title:    str = "Imported Course",
        language: str = "en",
        level:    str = "université",
        subject:  str = "",
    ) -> dict:
        log.info(f"📝 Construction depuis texte : {len(text)} chars")
        course_slug = self._course_slug(subject, title)
        return await self.structurer.structure(
            raw_text=text,
            language=language,
            level=level,
            title=title,
            subject=course_slug,
        )

    async def save_to_database(self, course_data: dict, db, domain: str = DEFAULT_DOMAIN) -> str:
        from database.models import Course, Chapter, Section, Concept
        log.info("💾 Sauvegarde PostgreSQL…")

        course = Course(
            title=course_data.get("title", "Cours importé"),
            domain=domain,  # 🎯 Stocker le domaine
            subject=course_data.get("subject", DEFAULT_COURSE),
            language=course_data.get("language", "en"),
            level=course_data.get("level", "université"),
            description=course_data.get("description", ""),
            file_path=course_data.get("file_path", ""),
        )
        db.add(course)
        await db.flush()

        slides = course_data.get("slides", [])  # 🎨 Récupérer les PNG paths

        for ch_data in course_data.get("chapters", []):
            chapter = Chapter(
                course_id=course.id,
                title=ch_data["title"],
                order=ch_data.get("order", 0),
                summary=ch_data.get("summary", ""),
            )
            db.add(chapter)
            await db.flush()

            for i, sec_data in enumerate(ch_data.get("sections", [])):
                # 🎨 Préserver le lien exact entre la section et sa slide PNG.
                section_order = sec_data.get("page_index") or sec_data.get("order") or (i + 1)
                image_url = (sec_data.get("image_url") or "").strip()

                if not image_url and section_order:
                    slide_index = int(section_order) - 1
                    if 0 <= slide_index < len(slides):
                        image_url = slides[slide_index]

                if not image_url and i < len(slides):
                    image_url = slides[i]

                image_urls = sec_data.get("image_urls") or ([] if not image_url else [image_url])
                
                section = Section(
                    chapter_id=chapter.id,
                    title=sec_data["title"],
                    order=section_order,
                    content=sec_data.get("content", ""),
                    image_url=image_url,  # 🎨 PNG path de cette slide
                    image_urls=image_urls,
                    duration_s=sec_data.get("duration_s", 120),
                )
                db.add(section)
                await db.flush()

                for c_data in sec_data.get("concepts", []):
                    concept = Concept(
                        section_id=section.id,
                        term=c_data.get("term", ""),
                        definition=c_data.get("definition", ""),
                        example=c_data.get("example", ""),
                        concept_type=c_data.get("type", "definition"),
                    )
                    db.add(concept)

        await db.commit()
        course_id = str(course.id)
        log.info(f"✅ Cours sauvegardé : ID={course_id}")
        return course_id

    def _estimate_duration(self, text: str) -> int:
        """Estime une durée de lecture (en secondes) selon le nombre de mots."""
        words = len((text or "").split())
        # ~130 mots/minute avec bornes raisonnables
        return max(30, min(600, int((words / 130.0) * 60)))

    def _infer_page_title(self, page_text: str, page_index: int) -> str:
        """Infer a page title by scoring heading-like lines.

        This avoids selecting formula annotations (e.g., "n: ...") as section titles.
        """
        lines = [ln.strip() for ln in (page_text or "").splitlines() if ln.strip()]
        if not lines:
            return f"Page {page_index}"

        candidates = lines[:12]
        best_line = ""
        best_score = -10.0

        for pos, line in enumerate(candidates):
            words = line.split()
            if len(line) < 3:
                continue
            if len(words) > 12:
                continue

            score = 0.0
            score += max(0.0, 4.0 - (pos * 0.45))

            if 2 <= len(words) <= 7:
                score += 2.0
            elif len(words) <= 10:
                score += 0.6

            alpha_ratio = sum(1 for ch in line if ch.isalpha()) / max(1, len(line))
            score += alpha_ratio * 2.0

            if line[:1].isupper():
                score += 0.8

            if ":" in line:
                score -= 1.7
            if re.search(r"\b(?:n|x|y|z|p|q|m|k|f)\s*:", line.lower()):
                score -= 2.2
            if re.search(r"[=<>±∑∫√]", line):
                score -= 2.0
            if re.search(r"\b\d+(?:\.\d+)?\b", line):
                score -= 0.6
            if line.endswith(":"):
                score -= 1.3

            if score > best_score:
                best_score = score
                best_line = line

        normalized = (best_line or "").strip(" .:-")
        if not normalized:
            normalized = lines[0][:80]

        return normalized[:100] if normalized else f"Page {page_index}"

    def _split_text_into_sections(self, raw_text: str) -> list[tuple[str, str]]:
        """
        Découpe un texte brut en sections à partir des lignes de titre probables.
        Retourne [(title, content), ...].
        """
        lines = [ln.rstrip() for ln in raw_text.splitlines()]
        cleaned = [ln for ln in lines if ln.strip()]
        if not cleaned:
            return []

        title_re = re.compile(
            r"^(?:chapter|chapitre|section|part|partie|module|lesson|lecon)\s*[0-9ivx\-:. ]+.*$",
            re.IGNORECASE,
        )
        numbered_re = re.compile(r"^[0-9]+(?:\.[0-9]+)*[\)\.]?\s+.+$")

        sections: list[tuple[str, list[str]]] = []
        current_title = "Introduction"
        current_lines: list[str] = []

        for ln in cleaned:
            is_short = len(ln) <= 120
            is_upper = ln.isupper() and len(ln) >= 4
            is_title = is_short and (title_re.match(ln) or numbered_re.match(ln) or is_upper)

            if is_title and current_lines:
                sections.append((current_title, current_lines))
                current_title = ln.strip()
                current_lines = []
            elif is_title and not current_lines and current_title == "Introduction":
                current_title = ln.strip()
            else:
                current_lines.append(ln)

        if current_lines:
            sections.append((current_title, current_lines))

        if not sections:
            text = "\n".join(cleaned)
            return [("Contenu", text)]

        out: list[tuple[str, str]] = []
        for title, body_lines in sections:
            content = "\n".join(body_lines).strip()
            if content:
                out.append((title.strip() or "Section", content))

        if not out:
            text = "\n".join(cleaned)
            return [("Contenu", text)]
        return out

    def _print_summary(self, data: dict) -> None:
        chapters = data.get("chapters", [])
        total_sections = sum(len(ch.get("sections", [])) for ch in chapters)
        total_concepts = sum(
            len(sec.get("concepts", []))
            for ch in chapters
            for sec in ch.get("sections", [])
        )
        total_duration = sum(
            sec.get("duration_s", 120)
            for ch in chapters
            for sec in ch.get("sections", [])
        )
        log.info(f"\n{'='*60}")
        log.info(f"✅ COURS : {data.get('title')}")
        log.info(f"   Matière    : {data.get('subject')}")
        log.info(f"   Niveau     : {data.get('level')}")
        log.info(f"   Langue     : {data.get('language')}")
        log.info(f"   Chapitres  : {len(chapters)}")
        log.info(f"   Sections   : {total_sections}")
        log.info(f"   Concepts   : {total_concepts}")
        log.info(f"   Durée est. : {total_duration // 60} min")
        log.info(f"{'='*60}\n")
        for i, ch in enumerate(chapters):
            log.info(f"  📖 Ch{i+1} : {ch['title']}")
            for j, sec in enumerate(ch.get("sections", [])):
                log.info(f"     └─ §{j+1} : {sec['title']} ({sec.get('duration_s',120)}s)")